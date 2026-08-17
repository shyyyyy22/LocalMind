from .base import BaseParser
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

class WordParser(BaseParser):
    def parse(self,path):
        doc = Document(path)
        content = "\n".join(para.text for para in doc.paragraphs)
        for table in doc.tables:
            for row in table.rows:
                content += " ".join(str(cell.text) for cell in row if cell.text)
                content += "\n"
        return content
class PresentationParser(BaseParser):
    def parse(self,path):
        ppt = Presentation(path)
        content = []
        for slide in ppt.slides:
            content.append("\n".join(shape.text for shape in slide.shapes if shape.has_text_frame))
        return '\n'.join(content)
class ExcelParser(BaseParser):
    def parse(self,path):
        workbook = load_workbook(path)
        content = ""
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows():
                content += " ".join(str(cell.value) for cell in row if cell.value)
                content += "\n"
        return content